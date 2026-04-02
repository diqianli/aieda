#!/usr/bin/env python3
"""
Agentic AI CPU 设计空间可视化 PPT 生成器
使用华为红色配色方案
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor

# 别名，保持代码一致性
RgbColor = RGBColor
from pptx.enum.shapes import MSO_SHAPE

# 华为红色配色方案
HUAWEI_RED = RgbColor(200, 16, 46)      # 主色
HUAWEI_DARK_RED = RgbColor(139, 0, 0)   # 深红
HUAWEI_LIGHT_RED = RgbColor(232, 93, 117) # 浅红
GRAY_BG = RgbColor(245, 245, 245)        # 浅灰背景
WHITE = RgbColor(255, 255, 255)
DARK_GRAY = RgbColor(51, 51, 51)
GOLD = RgbColor(245, 166, 35)            # 强调色

# 维度数据配置
DIMENSIONS = {
    1: {
        "title": "系统维度设计空间",
        "subtitle": "从机架级和多节点架构视角探索Agentic AI的系统级设计选择",
        "badge": "Level 1",
        "cards": [
            {
                "icon": "⚖️",
                "title": "CPU-GPU配比模式",
                "content": "决定CPU与GPU的数量比例，影响训练/推理效率和资源利用率。根据DeepSeek-V3优化经验，73%的性能提升来自CPU端优化。",
                "options": ["GPU为主型 (1:4)", "训练优化 (1:2)", "均衡型 (1:1)", "CPU机柜"]
            },
            {
                "icon": "🗄️",
                "title": "机架功率密度",
                "content": "优化功率密度、散热方式。Arm AGI CPU支持36kW风冷机架，8160核心。高密度部署降低TCO。",
                "options": ["标准 (20kW)", "高密度 (36kW)", "液冷 (60kW+)"]
            },
            {
                "icon": "🏗️",
                "title": "分层系统架构",
                "content": "将工作负载分配到不同层级。NVIDIA Vera Rubin支持独立CPU机柜（256颗CPU），专为RL/Agent推理设计。",
                "options": ["分离三层", "统一调度", "混合模式"]
            },
            {
                "icon": "🔗",
                "title": "节点间互联拓扑",
                "content": "决定多节点间的通信带宽和延迟。NVLink-C2C提供1.8TB/s一致性带宽，是Grace的2倍。",
                "options": ["NVLink-C2C (1.8TB/s)", "InfiniBand", "RoCEv2"]
            },
            {
                "icon": "💾",
                "title": "系统级内存架构",
                "content": "内存池化、分解式内存和CXL支持。Arm AGI CPU支持CXL 3.0 Type 3内存扩展。",
                "options": ["本地内存", "CXL 3.0池化", "分解式内存"]
            },
            {
                "icon": "🏭",
                "title": "部署模式",
                "content": "AI工厂 vs 传统数据中心。AI工厂强调机架级优化和大规模并行Agent环境。",
                "options": ["传统DC", "AI工厂", "边缘部署"]
            }
        ],
        "insight": "💡 技术洞察：阿姆达尔定律在AI系统中的应用——NVIDIA明确指出'系统性能越来越受制于agentic循环中CPU端串行任务的限制'。DeepSeek-V3在GB200上的优化显示，73%的约349 TFLOPS性能提升来自解决CPU overhead。"
    },
    2: {
        "title": "SoC维度设计空间",
        "subtitle": "探索芯片级架构设计，包括核心选择、缓存层次、互联和内存子系统",
        "badge": "Level 2",
        "cards": [
            {
                "icon": "🔬",
                "title": "核心架构选择",
                "content": "NVIDIA Vera采用自研Olympus核心（10-wide前端），Arm AGI CPU采用Neoverse V3。自研核心可实现1.5x IPC提升。",
                "options": ["公版核心", "自研核心", "混合架构"]
            },
            {
                "icon": "📚",
                "title": "缓存层次结构",
                "content": "Grace将L2从2MB缩减到1MB导致性能问题。Vera L3提升至162MB（+42%），Arm AGI CPU每核2MB L2。",
                "options": ["能效优先", "大缓存配置", "自适应缓存"]
            },
            {
                "icon": "🔌",
                "title": "Die架构策略",
                "content": "Vera采用单片架构消除NUMA边界，第二代SCF提供3.4TB/s分切带宽。Arm AGI CPU采用chiplet设计。",
                "options": ["单片Monolithic", "Chiplet", "混合架构"]
            },
            {
                "icon": "💾",
                "title": "内存子系统",
                "content": "Vera的SOCAMM支持1.5TB可更换内存（vs Grace焊死480GB）。Arm AGI CPU使用传统DDR5-8800。",
                "options": ["DDR5 DIMM", "SOCAMM", "HBM"]
            },
            {
                "icon": "🚪",
                "title": "I/O配置",
                "content": "Arm AGI CPU提供96条PCIe Gen6通道，Vera支持PCIe Gen6和1.8TB/s NVLink-C2C。",
                "options": ["标准配置", "高密度I/O", "定制I/O"]
            },
            {
                "icon": "🏭",
                "title": "工艺节点",
                "content": "Arm AGI CPU采用TSMC 3nm，Vera预计使用类似工艺。3nm相比5nm提供约1.7x密度提升和1.15x性能提升。",
                "options": ["5nm成熟", "3nm领先", "2nm前沿"]
            }
        ],
        "insight": "💡 技术洞察：单片 vs Chiplet的权衡——Vera的单片架构确保所有核心到资源的距离相同，'无需传统NUMA调优即可获得最佳性能'。SOCAMM的创新首次将LPDDR的能效优势与服务器级的可维护性结合。"
    },
    3: {
        "title": "CPU维度设计空间",
        "subtitle": "深入CPU核心微架构，探索多线程、指令集、分支预测和专用加速器",
        "badge": "Level 3",
        "cards": [
            {
                "icon": "🧵",
                "title": "多线程策略",
                "content": "Grace不支持SMT（72核=72线程）被证明是市场短板。Vera的空间多线程提供88核×2=176线程，物理隔离保证可预测延迟。",
                "options": ["无SMT", "传统SMT", "空间多线程"]
            },
            {
                "icon": "📝",
                "title": "AI指令集扩展",
                "content": "Vera是首个支持FP8的CPU。FP8可减少数据移动和转换开销，对于Agentic AI的频繁推理决策至关重要。",
                "options": ["BF16/INT8", "+FP8原生", "定制扩展"]
            },
            {
                "icon": "🔮",
                "title": "分支预测器",
                "content": "Python/PyTorch dispatch是分支密集型代码。Vera的神经分支预测器每周期可评估两个已采取分支。",
                "options": ["静态预测", "动态预测", "神经预测器"]
            },
            {
                "icon": "📏",
                "title": "指令前端宽度",
                "content": "Vera的10-wide前端宽于AMD Zen5 (8-wide)和Intel Granite Rapids (6-wide)，更宽前端提升ILP但增加功耗。",
                "options": ["6-wide (能效)", "8-wide (均衡)", "10-wide (性能)"]
            },
            {
                "icon": "🚀",
                "title": "片上专用加速器",
                "content": "针对AI工作负载的专用加速：RDMA引擎绕过CPU NoC，加密引擎降低安全开销，压缩引擎加速模型传输。",
                "options": ["最小配置", "标准加速", "全加速套件"]
            },
            {
                "icon": "⚡",
                "title": "频率与功耗策略",
                "content": "Agentic AI的kernel launch延迟敏感场景需要部分核心能boost到最高频率（DeepSeek文档建议）。",
                "options": ["能效优先", "平衡策略", "激进boost"]
            }
        ],
        "insight": "💡 技术洞察：空间多线程 vs 传统SMT——传统SMT通过时间切片共享流水线资源，导致性能抖动和尾延迟问题。空间多线程物理分区资源，保证每线程性能可预测，特别适合多租户AI工厂。"
    },
    4: {
        "title": "软硬件协同设计空间",
        "subtitle": "探索软硬件协同优化，从瓶颈识别到调度策略再到生态兼容",
        "badge": "Level 4",
        "cards": [
            {
                "icon": "🔍",
                "title": "CPU端瓶颈优化",
                "content": "DeepSeek-V3优化显示73%收益来自CPU端。Vera的1.5x IPC、FP8支持、神经预测器都针对此问题。",
                "options": ["基础优化", "激进优化", "软硬协同"]
            },
            {
                "icon": "🗺️",
                "title": "NUMA感知调度",
                "content": "DeepSeek使用bindpcie工具绑定进程到本地GPU获得70.6 TFLOPS提升。自动化NUMA调度是关键方向。",
                "options": ["手动绑定", "自动调度", "混合模式"]
            },
            {
                "icon": "🔧",
                "title": "软件栈优化",
                "content": "从Python调度优化到JIT编译，从PyTorch dispatch优化到kernel批处理，全栈协同消除CPU瓶颈。",
                "options": ["驱动层", "框架层", "全栈优化"]
            },
            {
                "icon": "🌐",
                "title": "生态兼容性策略",
                "content": "Grace初期遭遇TensorRT-LLM不支持arm64等问题。经过两年生态建设，Arm数据中心软件生态已大幅改善。",
                "options": ["ARM原生", "混合部署", "二进制翻译"]
            },
            {
                "icon": "📊",
                "title": "性能可观测性",
                "content": "统一的CPU-GPU协同性能分析工具，识别瓶颈并验证优化效果。关键是将CPU-GPU事件对齐到同一时间轴。",
                "options": ["分离工具", "统一视图", "AI分析"]
            },
            {
                "icon": "🎛️",
                "title": "运行时动态调优",
                "content": "根据工作负载特征动态调整SMT模式、核心频率和功耗封顶。Vera支持运行时切换单/双线程模式。",
                "options": ["静态配置", "动态调优", "预测性调优"]
            }
        ],
        "insight": "💡 技术洞察：CPU瓶颈的具体表现——在GB200上，Grace的kernel launch开销约为x86 Xeon的2倍。Blackwell GPU算力提升2.5倍后，CPU端'掩盖'GPU kernel间隙的时间急剧缩短，使瓶颈更加明显。"
    }
}


def add_red_header(slide, title_text, subtitle_text=""):
    """添加红色标题头"""
    # 标题区域背景
    header_shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        0, 0, Inches(10), Inches(1.3)
    )
    header_shape.fill.solid()
    header_shape.fill.fore_color.rgb = HUAWEI_RED
    header_shape.line.fill.background()

    # 标题文本
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.5))
    title_frame = title_box.text_frame
    title_frame.text = title_text
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(32)
    title_para.font.bold = True
    title_para.font.color.rgb = WHITE
    title_para.alignment = PP_ALIGN.LEFT

    # 副标题文本
    if subtitle_text:
        subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.8), Inches(9), Inches(0.4))
        subtitle_frame = subtitle_box.text_frame
        subtitle_frame.text = subtitle_text
        subtitle_para = subtitle_frame.paragraphs[0]
        subtitle_para.font.size = Pt(14)
        subtitle_para.font.color.rgb = RgbColor(255, 200, 200)
        subtitle_para.alignment = PP_ALIGN.LEFT


def add_card(slide, left, top, width, height, icon, title, content, options):
    """添加设计空间卡片"""
    # 卡片背景
    card_shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        left, top, width, height
    )
    card_shape.fill.solid()
    card_shape.fill.fore_color.rgb = WHITE
    card_shape.line.color.rgb = HUAWEI_RED
    card_shape.line.width = Pt(2)

    # 图标和标题
    header_box = slide.shapes.add_textbox(left + Inches(0.1), top + Inches(0.1), width - Inches(0.2), Inches(0.35))
    header_frame = header_box.text_frame
    header_frame.word_wrap = True

    # 图标行
    p = header_frame.paragraphs[0]
    p.text = f"{icon} {title}"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = HUAWEI_RED

    # 内容
    content_box = slide.shapes.add_textbox(left + Inches(0.1), top + Inches(0.45), width - Inches(0.2), height - Inches(0.55))
    content_frame = content_box.text_frame
    content_frame.word_wrap = True
    content_frame.text = content

    for paragraph in content_frame.paragraphs:
        paragraph.font.size = Pt(9)
        paragraph.font.color.rgb = DARK_GRAY
        paragraph.line_spacing = 1.2

    # 选项标签（底部小字）
    if options:
        options_text = " | ".join(options)
        options_box = slide.shapes.add_textbox(left + Inches(0.1), top + height - Inches(0.35), width - Inches(0.2), Inches(0.25))
        options_frame = options_box.text_frame
        options_frame.word_wrap = True
        options_frame.text = options_text
        for paragraph in options_frame.paragraphs:
            paragraph.font.size = Pt(8)
            paragraph.font.color.rgb = RgbColor(100, 100, 100)
            paragraph.font.italic = True


def add_insight_box(slide, insight_text):
    """添加底部技术洞察框"""
    # 洞察框背景
    insight_shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(0.5), Inches(6.5), Inches(9), Inches(0.8)
    )
    insight_shape.fill.solid()
    insight_shape.fill.fore_color.rgb = GRAY_BG
    insight_shape.line.color.rgb = GOLD
    insight_shape.line.width = Pt(2)

    # 洞察文本
    insight_box = slide.shapes.add_textbox(Inches(0.7), Inches(6.6), Inches(8.6), Inches(0.6))
    insight_frame = insight_box.text_frame
    insight_frame.word_wrap = True
    insight_frame.text = insight_text

    for paragraph in insight_frame.paragraphs:
        paragraph.font.size = Pt(10)
        paragraph.font.color.rgb = DARK_GRAY


def add_badge(slide, badge_text, left, top):
    """添加Level徽章"""
    badge_shape = slide.shapes.add_shape(
        MSO_SHAPE.OVAL,
        left, top, Inches(1.2), Inches(0.5)
    )
    badge_shape.fill.solid()
    badge_shape.fill.fore_color.rgb = GOLD
    badge_shape.line.fill.background()

    badge_box = slide.shapes.add_textbox(left, top + Inches(0.1), Inches(1.2), Inches(0.3))
    badge_frame = badge_box.text_frame
    badge_frame.text = badge_text
    p = badge_frame.paragraphs[0]
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER


def create_dimension_slide(prs, level_num, dimension_data):
    """创建维度幻灯片"""
    # 使用空白布局
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # 添加红色标题头
    add_red_header(slide, dimension_data["title"], dimension_data["subtitle"])

    # 添加Level徽章
    add_badge(slide, dimension_data["badge"], Inches(8.5), Inches(0.4))

    # 添加6个卡片（2行3列）
    card_width = Inches(2.9)
    card_height = Inches(1.8)
    margin = Inches(0.15)
    start_left = Inches(0.3)
    start_top = Inches(1.5)

    for i, card in enumerate(dimension_data["cards"]):
        row = i // 3
        col = i % 3
        left = start_left + col * (card_width + margin)
        top = start_top + row * (card_height + margin)

        add_card(
            slide,
            left, top, card_width, card_height,
            card["icon"],
            card["title"],
            card["content"],
            card["options"]
        )

    # 添加技术洞察框
    add_insight_box(slide, dimension_data["insight"])

    return slide


def create_ppt():
    """创建完整的PPT"""
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # 创建4页维度页
    for level in range(1, 5):
        create_dimension_slide(prs, level, DIMENSIONS[level])

    return prs


def main():
    """主函数"""
    print("正在生成 Agentic AI CPU 设计空间可视化 PPT...")
    print(f"配色方案: 华为红 (RGB: 200, 16, 46)")

    prs = create_ppt()

    output_path = "/Users/mac/storage/ptoa3/simpler_baseline/tauri-viz/Agentic_AI_CPU_Design_Space.pptx"
    prs.save(output_path)

    print(f"✅ PPT 已生成: {output_path}")
    print(f"📊 共 {len(prs.slides)} 页")
    print("\n页面概览:")
    for i, slide in enumerate(prs.slides, 1):
        print(f"  第{i}页: {DIMENSIONS[i]['title']}")


if __name__ == "__main__":
    main()
